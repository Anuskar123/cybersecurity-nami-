import os
import json
import PyPDF2
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import glob

def extract_text_from_pdf(filepath):
    try:
        text = ""
        with open(filepath, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_docx(filepath):
    try:
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"

def extract_text_from_pptx(filepath):
    try:
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def extract_text_from_html(filepath):
    try:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        return f"Error extracting HTML: {str(e)}"

def main():
    base_dir = r"c:\Users\DEII\Downloads\cybersecurity(nami)"
    
    # Collect all supported files
    supported_extensions = ['.pdf', '.docx', '.pptx', '.html', '.htm']
    all_files = []
    
    for root, dirs, files in os.walk(base_dir):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext in supported_extensions:
                all_files.append(os.path.join(root, file))
                
    print(f"Found {len(all_files)} supported files for extraction.")
    
    extracted_data = {}
    
    for filepath in all_files:
        filename = os.path.basename(filepath)
        rel_path = os.path.relpath(filepath, base_dir)
        ext = os.path.splitext(filepath)[1].lower()
        
        print(f"Processing: {rel_path}", flush=True)
        
        try:
            if ext == '.pdf':
                content = extract_text_from_pdf(filepath)
            elif ext == '.docx':
                content = extract_text_from_docx(filepath)
            elif ext == '.pptx':
                content = extract_text_from_pptx(filepath)
            elif ext in ['.html', '.htm']:
                content = extract_text_from_html(filepath)
            else:
                content = "Unsupported extension."
        except Exception as e:
            content = f"Error: {e}"
            print(f"Error processing {rel_path}: {e}", flush=True)
            
        # Basic cleanup
        if content:
            # remove excess newlines
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            cleaned_content = '\n'.join(lines)
            
            if cleaned_content:
                extracted_data[rel_path] = cleaned_content

    output_file = os.path.join(base_dir, 'class material', 'all_content_extracted.json')
    
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(extracted_data, f, indent=4, ensure_ascii=False)
        
    print(f"Extraction complete! Saved to {output_file}")
    print(f"Successfully extracted {len(extracted_data)} files.")

if __name__ == "__main__":
    main()
